from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Título e conteúdo
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title

    tf = body_placeholder.text_frame
    tf.clear()  # Limpa texto padrão

    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT


def main():
    prs = Presentation()

    # Slide 1
    slide_layout = prs.slide_layouts[0]  # Título
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Projeto de Análise Automatizada de Documentos Eletrônicos com IA"
    subtitle = slide.placeholders[1]
    subtitle.text = "Data: [Inserir data]\nApresentador: [Seu nome]"

    # Slide 2
    add_slide(prs, "Objetivo do Projeto", [
        "• Automatizar a análise do sistema de documentos eletrônicos",
        "• Detectar documentos duplicados",
        "• Identificar trâmites realizados incorretamente",
        "• Encontrar erros no cadastro dos documentos",
        "• Melhorar a eficiência e qualidade do processo documental"
    ])

    # Slide 3
    add_slide(prs, "Desafios Identificados", [
        "• Volume grande de documentos e trâmites",
        "• Dados textuais e estruturados heterogêneos",
        "• Dificuldade em identificar duplicidades e erros manualmente",
        "• Necessidade de integração com sistema legado"
    ])

    # Slide 4
    add_slide(prs, "Visão Geral da Solução", [
        "• Uso de Inteligência Artificial para análise textual e estruturada",
        "• Pipeline automatizada para pré-processamento, treino e avaliação",
        "• API para integração com sistema existente",
        "• Dashboard para monitoramento e feedback"
    ])

    # Slide 5
    add_slide(prs, "Etapas do Projeto", [
        "1. Entendimento e Coleta: Mapear dados, entender tipos de documentos e erros",
        "2. Pré-processamento: Limpeza, normalização e preparação dos dados",
        "3. Extração de Features: Transformar texto em vetores, preparar dados estruturados",
        "4. Treinamento de Modelos: Construir e validar modelos de classificação e detecção",
        "5. Validação e Ajustes: Avaliar desempenho, ajustar hiperparâmetros",
        "6. Implementação e Deploy: API para inferência, integração com sistema, dashboards"
    ])

    # Slide 6
    content = [
        "Sistema de Documentos <--> API de IA <--> Modelo de IA",
        "                                   |",
        "                                   v",
        "                            Banco de Resultados",
        "                                   |",
        "                                   v",
        "                              Dashboard",
        "",
        "• API recebe documentos e retorna análises",
        "• Modelo treinado para detectar duplicados, erros e anomalias",
        "• Banco armazena logs e resultados para auditoria",
        "• Dashboard para revisão e feedback"
    ]
    add_slide(prs, "Arquitetura do Sistema", content)

    # Slide 7
    add_slide(prs, "Tecnologias Utilizadas", [
        "• Python: linguagem principal",
        "• Pandas, scikit-learn, XGBoost: manipulação de dados e modelos",
        "• spaCy: pré-processamento de texto em português",
        "• transformers (Hugging Face): embeddings avançados (opcional)",
        "• FastAPI: API REST para integração",
        "• Uvicorn: servidor ASGI para API",
        "• Joblib: serialização de modelos e vetorizadores",
        "• FAISS: busca rápida por similaridade (duplicados)",
        "• Docker (futuro): containerização e deploy"
    ])

    # Slide 8
    add_slide(prs, "Pipeline Automatizada", [
        "• Carregamento e limpeza dos dados",
        "• Extração e transformação das features",
        "• Treinamento e validação do modelo",
        "• Armazenamento dos artefatos (modelos, vetorizadores)",
        "• Pode ser executada via script Python",
        "• Facilita re-treinamento e atualização periódica"
    ])

    # Slide 9
    add_slide(prs, "Exemplos de Uso da API", [
        "• Endpoint REST para análise rápida",
        "• Entrada: texto do documento + campos auxiliares",
        "• Saída: indicação de erro, duplicidade ou inconsistência",
        "• Pode ser integrada ao sistema legado para análises em tempo real"
    ])

    # Slide 10
    add_slide(prs, "Benefícios Esperados", [
        "• Redução de erros manuais e retrabalho",
        "• Processos mais ágeis e confiáveis",
        "• Melhoria na governança documental",
        "• Base para futuras automações e análises avançadas"
    ])

    # Slide 11
    add_slide(prs, "Próximos Passos", [
        "• Testes com dados reais do sistema",
        "• Ajustes finos e melhoria dos modelos",
        "• Desenvolvimento do dashboard para monitoramento",
        "• Implantação da API em ambiente de produção",
        "• Planejamento de feedback loop para aprendizado contínuo",
        "• Avaliação da containerização e orquestração (Docker, Airflow)"
    ])

    # Slide 12
    add_slide(prs, "Perguntas e Discussão", [
        "• Abrir espaço para dúvidas, sugestões e alinhamento técnico"
    ])

    # Salvar arquivo
    prs.save("Projeto_Analise_Documentos_IA.pptx")
    print("Arquivo 'Projeto_Analise_Documentos_IA.pptx' criado com sucesso!")


if __name__ == "__main__":
    main()